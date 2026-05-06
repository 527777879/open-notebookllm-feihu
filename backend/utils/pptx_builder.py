"""
PPTX Builder - 用於建立可編輯的 PPTX 檔案
"""
import os
import logging
import base64
import io
from typing import List, Dict, Any, Optional
from pathlib import Path

logger = logging.getLogger(__name__)

# 嘗試導入 python-pptx，如果沒有安裝則提供備援
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx 未安裝，PPTX 導出功能將不可用")

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False


class PPTXBuilder:
    """PPTX 建構器類別，用於從結構化內容建立可編輯的 PPTX 檔案"""

    # 標準簡報尺寸 (16:9 比例)
    DEFAULT_SLIDE_WIDTH_INCHES = 13.333
    DEFAULT_SLIDE_HEIGHT_INCHES = 7.5

    # 預設 DPI
    DEFAULT_DPI = 96

    def __init__(self, slide_width_inches: float = None, slide_height_inches: float = None):
        """
        初始化 PPTX 建構器

        Args:
            slide_width_inches: 簡報寬度（英吋）
            slide_height_inches: 簡報高度（英吋）
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx 未安裝，請執行 pip install python-pptx")

        self.slide_width_inches = slide_width_inches or self.DEFAULT_SLIDE_WIDTH_INCHES
        self.slide_height_inches = slide_height_inches or self.DEFAULT_SLIDE_HEIGHT_INCHES
        self.prs = None
        self.current_slide = None

    def create_presentation(self) -> 'Presentation':
        """建立新的簡報"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.slide_width_inches)
        self.prs.slide_height = Inches(self.slide_height_inches)
        return self.prs

    def add_blank_slide(self):
        """新增空白簡報頁"""
        if not self.prs:
            self.create_presentation()

        # 使用空白版面配置（通常是 layout 6）
        blank_layout = self.prs.slide_layouts[6]
        self.current_slide = self.prs.slides.add_slide(blank_layout)
        return self.current_slide

    def add_title_slide(self, title: str, subtitle: str = None):
        """
        新增標題頁

        Args:
            title: 主標題
            subtitle: 副標題
        """
        if not self.prs:
            self.create_presentation()

        # 使用標題版面配置
        title_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(title_layout)

        # 設定標題
        if slide.shapes.title:
            slide.shapes.title.text = title

        # 設定副標題
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

        self.current_slide = slide
        return slide

    def add_content_slide(self, title: str, content: str, bullet_points: List[str] = None):
        """
        新增內容頁

        Args:
            title: 頁面標題
            content: 內容文字
            bullet_points: 條列要點
        """
        if not self.prs:
            self.create_presentation()

        # 使用標題和內容版面配置
        content_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(content_layout)

        # 設定標題
        if slide.shapes.title:
            slide.shapes.title.text = title

        # 設定內容
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame

            if bullet_points:
                for i, point in enumerate(bullet_points):
                    if i == 0:
                        tf.text = point
                    else:
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0
            elif content:
                tf.text = content

        self.current_slide = slide
        return slide

    def add_image_to_slide(self, slide, image_data: str,
                           left: float = 0, top: float = 0,
                           width: float = None, height: float = None):
        """
        將圖片加入簡報頁

        Args:
            slide: 目標簡報頁
            image_data: Base64 編碼的圖片資料
            left: 左邊距（英吋）
            top: 上邊距（英吋）
            width: 寬度（英吋）
            height: 高度（英吋）
        """
        if not PIL_AVAILABLE:
            logger.warning("PIL 未安裝，無法處理圖片")
            return

        try:
            # 解碼 Base64 圖片
            if image_data.startswith('data:image'):
                # 移除 data URL 前綴
                image_data = image_data.split(',')[1]

            image_bytes = base64.b64decode(image_data)
            image_stream = io.BytesIO(image_bytes)

            # 如果沒有指定尺寸，使用全頁
            if width is None:
                width = self.slide_width_inches
            if height is None:
                height = self.slide_height_inches

            # 新增圖片
            slide.shapes.add_picture(
                image_stream,
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(height)
            )
            logger.debug(f"圖片已加入簡報頁")

        except Exception as e:
            logger.error(f"加入圖片失敗: {e}")

    def add_text_box(self, slide, text: str,
                     left: float, top: float,
                     width: float, height: float,
                     font_size: int = 18,
                     bold: bool = False,
                     align: str = 'left'):
        """
        新增文字方塊

        Args:
            slide: 目標簡報頁
            text: 文字內容
            left, top, width, height: 位置和尺寸（英吋）
            font_size: 字體大小（點）
            bold: 是否粗體
            align: 對齊方式 ('left', 'center', 'right')
        """
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        tf = textbox.text_frame
        tf.text = text
        tf.word_wrap = True

        # 設定段落格式
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = bold

            if align == 'center':
                paragraph.alignment = PP_ALIGN.CENTER
            elif align == 'right':
                paragraph.alignment = PP_ALIGN.RIGHT
            else:
                paragraph.alignment = PP_ALIGN.LEFT

        return textbox

    def build_from_presentation_data(self, presentation_data: Dict[str, Any]) -> 'Presentation':
        """
        從簡報資料建立 PPTX

        Args:
            presentation_data: 簡報資料字典，包含 title, slides 等

        Returns:
            Presentation 物件
        """
        self.create_presentation()

        slides = presentation_data.get('slides', [])

        for slide_data in slides:
            slide_type = slide_data.get('type', 'content')
            title = slide_data.get('title', '')
            content = slide_data.get('content', '')
            image = slide_data.get('image')

            if slide_type == 'title':
                # 標題頁
                slide = self.add_blank_slide()

                # 如果有背景圖片，先加入
                if image:
                    self.add_image_to_slide(slide, image)

                # 加入標題文字
                self.add_text_box(
                    slide, title,
                    left=0.5, top=2.5,
                    width=12.333, height=1.5,
                    font_size=44, bold=True, align='center'
                )

                # 加入副標題
                if content:
                    self.add_text_box(
                        slide, content,
                        left=0.5, top=4.2,
                        width=12.333, height=1,
                        font_size=24, align='center'
                    )
            else:
                # 內容頁
                slide = self.add_blank_slide()

                # 如果有背景圖片
                if image:
                    self.add_image_to_slide(slide, image)

                # 加入標題
                self.add_text_box(
                    slide, title,
                    left=0.5, top=0.3,
                    width=12.333, height=0.8,
                    font_size=32, bold=True
                )

                # 加入內容
                if content:
                    # 解析條列內容
                    lines = content.split('\n')
                    formatted_content = '\n'.join(lines)

                    self.add_text_box(
                        slide, formatted_content,
                        left=0.5, top=1.3,
                        width=12.333, height=5.5,
                        font_size=20
                    )

        return self.prs

    def save(self, output_path: str):
        """
        儲存簡報到檔案

        Args:
            output_path: 輸出檔案路徑
        """
        if not self.prs:
            raise ValueError("沒有可儲存的簡報")

        # 確保目錄存在
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True)

        self.prs.save(output_path)
        logger.info(f"簡報已儲存至: {output_path}")

    def save_to_bytes(self) -> bytes:
        """
        將簡報儲存為 bytes

        Returns:
            PPTX 檔案的 bytes 資料
        """
        if not self.prs:
            raise ValueError("沒有可儲存的簡報")

        buffer = io.BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def save_to_base64(self) -> str:
        """
        將簡報儲存為 Base64 編碼

        Returns:
            Base64 編碼的 PPTX 資料
        """
        pptx_bytes = self.save_to_bytes()
        return base64.b64encode(pptx_bytes).decode('utf-8')


def is_pptx_available() -> bool:
    """檢查 PPTX 功能是否可用"""
    return PPTX_AVAILABLE
